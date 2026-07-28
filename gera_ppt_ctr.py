#!/usr/bin/env python3
"""Gera o PPT de qualidade de CTR da Vanzolini a partir dos dados do Supabase.

Os numeros nunca sao digitados aqui: o script busca tudo pelas RPCs e monta o
arquivo. Para atualizar o deck depois de um novo mes de midia, basta rodar de
novo. Uso: python3 gera_ppt_ctr.py <caminho-do-pptx>
"""
import json, subprocess, sys, urllib.request

URL = "https://ltasijrhkotyyrxnavab.supabase.co/rest/v1/rpc"
KEY = ("eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6Imx0YXNpanJoa290eXl"
       "yeG5hdmFiIiwicm9sZSI6ImFub24iLCJpYXQiOjE3ODExMjk5NDAsImV4cCI6MjA5NjcwNTk0MH0"
       ".XAJmbTSm6d5Y6xobOLceHlVr0e_iratHW_u6atzUC5c")
INI, FIM = "2026-01-01", "2026-06-30"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x06, 0x06, 0x53)
AZUL   = RGBColor(0x48, 0x48, 0xB5)
TXT    = RGBColor(0x1A, 0x1A, 0x2E)
MUT    = RGBColor(0x6F, 0x6F, 0x8A)
BORDA  = RGBColor(0xE4, 0xE4, 0xF1)
OFF    = RGBColor(0xF7, 0xF7, 0xFC)
VERDE  = RGBColor(0x1F, 0x8A, 0x4C)
VERDEBG= RGBColor(0xE8, 0xF5, 0xEE)
AMBBG  = RGBColor(0xFD, 0xF3, 0xE0)
AMB    = RGBColor(0x7A, 0x5B, 0x00)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA  = RGBColor(0xC3, 0xC3, 0xD4)
H1F, BF = "Inter Tight", "Work Sans"

# Referencias de mercado. Cada uma guarda a metodologia declarada pela fonte,
# porque e isso que permite comparar do mesmo jeito dos dois lados.
BENCH = {
    "Google":  dict(v=7.56, fonte="LocaliQ, Google Ads Benchmarks, setor Educação e Instrução, 2025",
                    met="média do setor na Rede de Pesquisa"),
    "Meta":    dict(v=1.80, fonte="Superads, benchmark de anúncios de Educação, 2025",
                    met="mediana entre campanhas"),
    "LinkedIn":dict(v=0.61, fonte="Benchmarks públicos de LinkedIn Ads, média geral de todos os setores",
                    met="média geral"),
}
MES = ["jan","fev","mar","abr","mai","jun","jul","ago","set","out","nov","dez"]


def rpc(nome, **args):
    req = urllib.request.Request(f"{URL}/{nome}", method="POST",
        data=json.dumps(args).encode(),
        headers={"apikey": KEY, "Authorization": f"Bearer {KEY}", "Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=60) as r:
        return json.loads(r.read())


def num(v, casas=0):
    s = f"{float(v or 0):,.{casas}f}"
    return s.replace(",", "@").replace(".", ",").replace("@", ".")

def brl(v, casas=0):  return "R$ " + num(v, casas)
def pct(v):           return "não rodou" if v is None else num(v, 2) + "%"
def rot(m):
    a, x = m.split("-"); return f"{MES[int(x)-1]}/{a[2:]}"


# ---------------------------------------------------------------- primitivas
def txbox(sl, x, y, w, h, texto, tam=12, cor=TXT, negrito=False, fonte=BF,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, espaco=1.0):
    cx = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = cx.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = espaco
    r = p.add_run(); r.text = texto
    r.font.size, r.font.bold, r.font.name, r.font.color.rgb = Pt(tam), negrito, fonte, cor
    return cx

def retangulo(sl, x, y, w, h, preench=None, linha=None, raio=0.06):
    forma = MSO_SHAPE.ROUNDED_RECTANGLE if raio else MSO_SHAPE.RECTANGLE
    s = sl.shapes.add_shape(forma, Inches(x), Inches(y), Inches(w), Inches(h))
    if raio:
        try: s.adjustments[0] = raio
        except Exception: pass
    if preench: s.fill.solid(); s.fill.fore_color.rgb = preench
    else:       s.fill.background()
    if linha: s.line.color.rgb = linha; s.line.width = Pt(0.75)
    else:     s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s

def linha_h(sl, x, y, w, cor=BORDA, esp=0.75):
    from pptx.util import Emu
    c = sl.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    c.line.color.rgb = cor; c.line.width = Pt(esp)
    return c


def cabecalho(sl, titulo, sub, tag, logo):
    if logo:
        sl.shapes.add_picture(logo, Inches(0.55), Inches(0.36), height=Inches(0.34))
    txbox(sl, 1.95, 0.31, 8.4, 0.34, titulo, 21, NAVY, True, H1F)
    txbox(sl, 1.95, 0.66, 8.6, 0.26, sub, 11, MUT)
    txbox(sl, 10.2, 0.33, 2.6, 0.24, tag, 13, NAVY, True, H1F, PP_ALIGN.RIGHT)
    txbox(sl, 10.2, 0.60, 2.6, 0.22, "janeiro a junho de 2026", 10, MUT, align=PP_ALIGN.RIGHT)
    linha_h(sl, 0.55, 1.02, 12.23, NAVY, 1.5)

def rodape(sl, txt):
    linha_h(sl, 0.55, 6.98, 12.23, BORDA)
    txbox(sl, 0.55, 7.06, 5.0, 0.22, "Fundação Vanzolini · Qualidade de mídia", 8.5, MUT)
    txbox(sl, 5.6, 7.06, 7.18, 0.32, txt, 8.5, MUT, align=PP_ALIGN.RIGHT)


def card(sl, x, y, w, h, plataforma, escopo, nosso, contexto):
    """Card de destaque: o CTR da Vanzolini lado a lado com a referencia da plataforma.

    `nosso` pode ser None, para plataforma que nao rodou no periodo. Nesse caso o
    card continua na tela, porque canal nao usado tambem e informacao.
    """
    b = BENCH[plataforma]
    retangulo(sl, x, y, w, h, BRANCO, BORDA)
    retangulo(sl, x, y, w, 0.07, AZUL, None, 0)
    txbox(sl, x+0.26, y+0.2, w-0.5, 0.3, plataforma, 17, NAVY, True, H1F)
    txbox(sl, x+0.26, y+0.52, w-0.5, 0.22, escopo, 10, MUT)

    meio = w * 0.52
    txbox(sl, x+0.26, y+0.82, meio-0.35, 0.2, "VANZOLINI", 8, MUT, True)
    txbox(sl, x+0.26, y+1.0, meio-0.35, 0.5,
          pct(nosso) if nosso is not None else "não rodou",
          30 if nosso is not None else 19, NAVY if nosso is not None else CINZA, True, H1F)
    # Caixa propria para a referencia, para o numero de mercado ter o mesmo peso visual.
    retangulo(sl, x+meio, y+0.76, w-meio-0.26, 0.86, OFF, BORDA, 0.08)
    txbox(sl, x+meio+0.18, y+0.86, w-meio-0.6, 0.2, "MERCADO", 8, MUT, True)
    txbox(sl, x+meio+0.18, y+1.04, w-meio-0.6, 0.5, pct(b["v"]), 30, AZUL, True, H1F)

    if nosso is None:
        rotulo, bom = "canal não testado", None
    else:
        # Ate 10% de diferenca a comparacao e ruido entre metodologias, nao desempenho.
        razao = nosso / b["v"]
        if razao >= 1.0:    rotulo, bom = f"{num(razao,1)}x acima do mercado", True
        elif razao >= 0.90: rotulo, bom = "em linha com o mercado", True
        else:               rotulo, bom = "abaixo da referência geral", False
    fundo = {True: VERDEBG, False: AMBBG, None: OFF}[bom]
    cor   = {True: VERDE,   False: AMB,   None: MUT}[bom]
    largura = 0.108 * len(rotulo) + 0.4
    retangulo(sl, x+0.26, y+1.7, largura, 0.3, fundo, None, 0.5)
    txbox(sl, x+0.26, y+1.755, largura, 0.22, rotulo, 10.5, cor, True, align=PP_ALIGN.CENTER)
    txbox(sl, x+0.26, y+2.1, w-0.52, 0.68, contexto, 9.5, MUT, espaco=1.2)


def tabela(sl, x, y, w, colunas, linhas, larguras, alt_lin=0.245, tam=10.5):
    """Tabela desenhada a mao: da controle fino sobre altura de linha e cor."""
    xs, acum = [], x
    for f in larguras:
        xs.append(acum); acum += f * w
    # cabecalho
    retangulo(sl, x, y, w, 0.28, NAVY, None, 0)
    for i, c in enumerate(colunas):
        al = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        txbox(sl, xs[i]+0.08, y+0.075, larguras[i]*w-0.16, 0.2, c.upper(), 8.5, BRANCO, True, align=al)
    yy = y + 0.28
    for n, ln in enumerate(linhas):
        estilo = ln.get("estilo")
        fundo = {"total": RGBColor(0xEE,0xF0,0xF8), "bench": VERDEBG}.get(estilo,
                 OFF if n % 2 else BRANCO)
        retangulo(sl, x, yy, w, alt_lin, fundo, None, 0)
        if estilo == "total": linha_h(sl, x, yy, w, NAVY, 1.25)
        for i, cel in enumerate(ln["cels"]):
            if isinstance(cel, tuple): valor, cor, neg = cel
            else:                      valor, cor, neg = cel, TXT, estilo in ("total","bench")
            if estilo == "bench": cor = RGBColor(0x15,0x5E,0x35)
            al = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            txbox(sl, xs[i]+0.08, yy+0.045, larguras[i]*w-0.16, alt_lin-0.05,
                  valor, tam, cor, neg, align=al)
        yy += alt_lin
    return yy


def linhas_evolucao(ev, colunas_plat):
    out, tl, ti = [], 0, 0
    med = lambda ks: sorted(v for v in ks if v is not None)
    for l in ev:
        tl += l["leads"]; ti += float(l["investimento"])
        cels = [rot(l["mes"]), num(l["leads"]), brl(l["investimento"]),
                brl(l["cpl"], 2) if l["cpl"] is not None else "sem dado"]
        for p in colunas_plat:
            v = l["ctr_" + p.lower()]
            cels.append((pct(v), CINZA if v is None else (VERDE if v >= BENCH[p]["v"] else TXT),
                         v is not None and v >= BENCH[p]["v"]))
        out.append(dict(cels=cels))
    cels = ["Semestre", num(tl), brl(ti), brl(ti/tl, 2) if tl else "sem dado"]
    for p in colunas_plat:
        vs = med([l["ctr_"+p.lower()] for l in ev])
        cels.append(pct(vs[len(vs)//2] if len(vs) % 2 else (vs[len(vs)//2-1]+vs[len(vs)//2])/2) if vs else "não rodou")
    out.append(dict(cels=cels, estilo="total"))
    out.append(dict(cels=["Referência de mercado", "", "", ""] + [pct(BENCH[p]["v"]) for p in colunas_plat],
                    estilo="bench"))
    return out


def bloco_campanhas(sl, x, y, w, plataforma, itens, formato):
    txbox(sl, x, y, w, 0.24, plataforma, 13, NAVY, True, H1F)
    linha_h(sl, x, y+0.26, w, BORDA)
    yy = y + 0.34
    for grupo, rotulo in (("CTR", "Maior CTR"), ("ESCALA", "Maior escala")):
        g = [i for i in itens if i["destaque"] == grupo]
        if not g: continue
        txbox(sl, x, yy, w, 0.18, rotulo.upper(), 8, MUT, True)
        yy += 0.2
        for it in g:
            # Nome longo quebraria em duas linhas e cobriria o formato logo abaixo.
            nome = it["curso"]
            if len(nome) > 52: nome = nome[:51].rstrip(" ,") + "…"
            txbox(sl, x, yy, w-0.85, 0.2, nome, 10, TXT)
            txbox(sl, x, yy+0.185, w-0.85, 0.18, formato(it["campanha"]), 8, MUT)
            txbox(sl, x+w-0.85, yy+0.02, 0.85, 0.2, pct(it["ctr"]), 11, NAVY, True, align=PP_ALIGN.RIGHT)
            yy += 0.37
        yy += 0.04
    return yy


def formato(c):
    n = (c or "").lower()
    if n.startswith("search"):  return "Rede de Pesquisa"
    if n.startswith("display"): return "Display e Demand Gen"
    if "lkd" in n or "linkedin" in n: return "LinkedIn, geração de leads"
    return "Meta, geração de leads"


# ------------------------------------------------------------------- montagem
def main(destino, logo):
    d = {
        "cons_mba":  rpc("consolidado_plataforma", p_ini=INI, p_fim=FIM, p_tipo="MBA"),
        "cons_curso":rpc("consolidado_plataforma", p_ini=INI, p_fim=FIM, p_tipo="Curso"),
        "gsd":       rpc("google_search_display", p_ini=INI, p_fim=FIM),
        "ev_mba":    rpc("evolucao_mensal_ctr", p_ini=INI, p_fim=FIM, p_tipo="MBA"),
        "ev_curso":  rpc("evolucao_mensal_ctr", p_ini=INI, p_fim=FIM, p_tipo="Curso"),
    }
    for t in ("MBA", "Curso"):
        for p in ("Google", "Meta", "LinkedIn"):
            d[f"ex_{t}_{p}"] = rpc("campanhas_exemplo_tipo", p_ini=INI, p_fim=FIM,
                                   p_plataforma=p, p_tipo=t)
    acha = lambda lista, p: next((x for x in lista if x["plataforma"] == p), None)
    pesquisa = next(g for g in d["gsd"] if g["formato"] == "Rede de Pesquisa")
    display  = next(g for g in d["gsd"] if g["formato"] == "Display e Demand Gen")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    branco = prs.slide_layouts[6]
    RODAPE = "Método e fontes na página de anotações"
    COLS = [0.55, 4.685, 8.82]   # tres cards, 3,91 de largura, ocupando a faixa util
    LARG, ALTO, TOPO = 3.91, 2.86, 1.2
    COLUNAS = ["Mês", "Leads", "Investimento", "CPL", "CTR Google", "CTR Meta", "CTR LinkedIn"]
    LARGURAS = [0.16, 0.13, 0.17, 0.13, 0.14, 0.135, 0.135]

    # ---- Slide 1: Cursos de curta duração
    sl = prs.slides.add_slide(branco)
    cabecalho(sl, "Cursos de curta duração",
              "CTR das campanhas contra a referência de mercado, e a evolução mês a mês",
              "Primeiro semestre", logo)
    m, li = acha(d["cons_curso"], "Meta"), acha(d["cons_curso"], "LinkedIn")
    card(sl, COLS[0], TOPO, LARG, ALTO, "Google", "não rodou no semestre", None,
         "Nenhum curso usou Google no período. É o canal com maior espaço para teste "
         "nesta linha.")
    card(sl, COLS[1], TOPO, LARG, ALTO, "Meta", f"{m['campanhas']} campanhas", m["ctr_mediano"],
         f"O CTR subiu de {pct(d['ev_curso'][0]['ctr_meta'])} em janeiro para "
         f"{pct(d['ev_curso'][-1]['ctr_meta'])} em junho, com o CPL caindo no mesmo período.")
    card(sl, COLS[2], TOPO, LARG, ALTO, "LinkedIn", f"{li['campanhas']} campanhas", li["ctr_mediano"],
         f"{num(li['impressoes'])} impressões no semestre. É o canal que fica acima da "
         f"referência com mais consistência.")
    tabela(sl, 0.55, 4.26, 12.23, COLUNAS,
           linhas_evolucao(d["ev_curso"], ["Google", "Meta", "LinkedIn"]), LARGURAS)
    rodape(sl, RODAPE)

    # ---- Slide 2: MBAs
    sl = prs.slides.add_slide(branco)
    cabecalho(sl, "MBAs", "CTR das campanhas contra a referência de mercado, e a evolução mês a mês",
              "Primeiro semestre", logo)
    m, li = acha(d["cons_mba"], "Meta"), acha(d["cons_mba"], "LinkedIn")
    card(sl, COLS[0], TOPO, LARG, ALTO, "Google", "Rede de Pesquisa", pesquisa["ctr_mediano"],
         f"{pesquisa['campanhas']} campanhas de Pesquisa. Nas de Display e Demand Gen o CTR "
         f"mediano foi de {pct(display['ctr_mediano'])}.")
    card(sl, COLS[1], TOPO, LARG, ALTO, "Meta", f"{m['campanhas']} campanhas", m["ctr_mediano"],
         f"As campanhas de maior alcance do semestre, {num(m['impressoes'])} impressões. "
         f"Alcance amplo derruba o CTR e sustenta o volume.")
    card(sl, COLS[2], TOPO, LARG, ALTO, "LinkedIn", f"{li['campanhas']} campanhas", li["ctr_mediano"],
         f"Só {num(li['impressoes'])} impressões no semestre. O canal rodou pouco nos MBAs, "
         f"então o número ainda é frágil.")
    tabela(sl, 0.55, 4.26, 12.23, COLUNAS,
           linhas_evolucao(d["ev_mba"], ["Google", "Meta", "LinkedIn"]), LARGURAS)
    rodape(sl, RODAPE)

    # ---- Slide 3: campanhas por plataforma, MBA em cima e Cursos embaixo
    sl = prs.slides.add_slide(branco)
    cabecalho(sl, "Onde o CTR se sustenta", "As campanhas de maior CTR e as de maior escala em cada plataforma",
              "Por plataforma", logo)
    txbox(sl, 0.55, 1.16, 4.0, 0.26, "MBAs", 15, NAVY, True, H1F)
    linha_h(sl, 0.55, 1.44, 12.23, AZUL, 1.25)
    bloco_campanhas(sl, 0.55, 1.56, 5.95, "Google", d["ex_MBA_Google"], formato)
    bloco_campanhas(sl, 6.83, 1.56, 5.95, "Meta", d["ex_MBA_Meta"], formato)
    txbox(sl, 0.55, 4.06, 4.0, 0.26, "Cursos de curta duração", 15, NAVY, True, H1F)
    linha_h(sl, 0.55, 4.34, 12.23, AZUL, 1.25)
    bloco_campanhas(sl, 0.55, 4.46, 5.95, "Meta", d["ex_Curso_Meta"], formato)
    bloco_campanhas(sl, 6.83, 4.46, 5.95, "LinkedIn", d["ex_Curso_LinkedIn"], formato)
    li_mba = acha(d["cons_mba"], "LinkedIn")
    linha_h(sl, 0.55, 6.98, 12.23, BORDA)
    txbox(sl, 0.55, 7.06, 8.4, 0.3,
          f"O LinkedIn não aparece nos exemplos de MBA porque rodou pouco no semestre, "
          f"{num(li_mba['impressoes'])} impressões em {li_mba['campanhas']} campanhas.", 8.5, MUT)
    txbox(sl, 9.1, 7.06, 3.68, 0.3, "Campanhas com pelo menos 20 mil impressões no semestre",
          8.5, MUT, align=PP_ALIGN.RIGHT)

    # ---- Slide 4: metodologia, para quem vai apresentar
    sl = prs.slides.add_slide(branco)
    cabecalho(sl, "Como estes números foram calculados",
              "Página de apoio para quem apresenta, não precisa entrar na versão do cliente",
              "Metodologia", logo)
    blocos = [
        ("Mediana, não média",
         "O número da Vanzolini é a mediana do CTR entre as campanhas. A média seria puxada por uma "
         "campanha isolada de volume muito alto ou muito baixo; a mediana mostra o desempenho da campanha típica."),
        ("O que ficou de fora",
         "Campanhas institucionais, que são de alcance e não de clique, e campanhas com menos de mil "
         "impressões no semestre, que ainda não têm volume para gerar um CTR estável."),
        ("Por que o Google aparece só como Rede de Pesquisa",
         "A referência da LocaliQ é de Rede de Pesquisa. Comparar a mediana geral do Google, que inclui "
         "Display e Demand Gen, contra um benchmark de Pesquisa inflaria o resultado a nosso favor. "
         "O Display aparece à parte, sem comparação, porque não há referência equivalente."),
        ("Cada plataforma com a referência dela",
         "CTR não é comparável entre plataformas. No Google a pessoa buscou pelo nome do curso, então o "
         "CTR é naturalmente alto. No Meta fica entre 1% e 2%, e no LinkedIn abaixo de 1%. "
         "A leitura correta é sempre dentro da mesma coluna."),
        ("Maior CTR e maior escala não são a mesma coisa",
         "As campanhas de maior alcance têm CTR menor porque falam com um público mais amplo, que ainda "
         "não conhece o curso. É o caso dos MBAs no Meta, e é esse volume que sustenta o funil. "
         "Nenhuma das duas leituras sozinha conta a história."),
        ("Quando dizemos que está em linha",
         "Diferença de até 10% para a referência é tratada como em linha, e não como abaixo. "
         "A variação entre metodologias de benchmark é maior que isso, então apontar uma diferença "
         "de 8% como queda seria criar um problema que os dados não sustentam."),
    ]
    for i, (t, c) in enumerate(blocos):
        col, lin = i % 2, i // 2
        x, y = 0.55 + col * 6.28, 1.2 + lin * 1.72
        retangulo(sl, x, y, 5.95, 1.6, OFF, BORDA)
        retangulo(sl, x, y, 0.055, 1.6, AZUL, None, 0)
        txbox(sl, x+0.28, y+0.16, 5.4, 0.26, t, 12, NAVY, True, H1F)
        txbox(sl, x+0.28, y+0.48, 5.4, 1.0, c, 9, MUT, espaco=1.28)
    y = 1.2 + 3 * 1.72
    retangulo(sl, 0.55, y, 12.23, 0.62, BRANCO, BORDA)
    txbox(sl, 0.83, y+0.11, 11.6, 0.2, "REFERÊNCIAS CITADAS", 8, MUT, True)
    txbox(sl, 0.83, y+0.31, 11.6, 0.24,
          " · ".join(f"{p}: {b['fonte']} ({b['met']})" for p, b in BENCH.items()), 8.5, TXT)
    rodape(sl, "Dados: plataformas de mídia e RD Station · consolidado pela Communitas")

    prs.save(destino)
    print("ok:", destino)


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
